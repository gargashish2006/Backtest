from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)  # Widescreen 16:9
prs.slide_height = Inches(7.5)

# Define colors
bg_color = RGBColor(15, 23, 42)    # Dark Slate (Background)
accent_color = RGBColor(16, 185, 129)  # Emerald Green (Accent)
text_color = RGBColor(248, 250, 252)   # White/Light Gray (Main Text)
mute_color = RGBColor(148, 163, 184)   # Slate Gray (Secondary Text)
darker_bg = RGBColor(30, 41, 59)       # Slightly lighter bg for shapes
highlight_color = RGBColor(255, 172, 65) # Orange for emphasis

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color

def add_footer(slide, page_num):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.0), Inches(13.333), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = darker_bg
    shape.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(6), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "CONFIDENTIAL | Quality Growth Value Fundamental Strategy"
    p.font.size = Pt(10)
    p.font.color.rgb = mute_color
    
    txBox = slide.shapes.add_textbox(Inches(12.3), Inches(7.1), Inches(1), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(page_num)
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(10)
    p.font.color.rgb = mute_color

def add_master_layout_elements(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.2), Inches(0.15), Inches(0.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()

def create_title_slide(prs, title, subtitle, tagline):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3), Inches(0.4), Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = accent_color
    shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.font.name = 'Arial'
    
    p = tf.add_paragraph()
    p.text = subtitle
    p.font.size = Pt(32)
    p.font.color.rgb = accent_color
    p.font.name = 'Arial'
    p.space_before = Pt(10)

    tag_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.33), Inches(1))
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.text = tagline
    p.font.size = Pt(20)
    p.font.color.rgb = mute_color
    p.font.name = 'Arial'

    return slide

def create_content_slide(prs, title, subtitle, bullets, page=0, disclaimer=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_master_layout_elements(slide)
    add_footer(slide, page)
    
    t_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    tf = t_box.text_frame
    p = tf.paragraphs[0]
    p.text = title.upper()
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.font.name = 'Arial'
    
    p = tf.add_paragraph()
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = accent_color
    p.font.name = 'Arial'

    text_width = Inches(12)
    b_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), text_width, Inches(4.5))
    tf = b_box.text_frame
    tf.word_wrap = True
    
    for item in bullets:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(22)
        p.font.color.rgb = text_color
        p.font.name = 'Arial'
        p.space_after = Pt(20)
        p.space_before = Pt(5)
        p.level = 0
    
    if disclaimer:
        d_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.5))
        tf = d_box.text_frame
        p = tf.paragraphs[0]
        p.text = disclaimer
        p.font.size = Pt(12)
        p.font.color.rgb = mute_color
        p.font.italic = True

    return slide

def create_emphasis_slide(prs, title, main_text, highlight_text, page=0):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_footer(slide, page)
    
    t_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.33), Inches(1.5))
    tf = t_box.text_frame
    p = tf.paragraphs[0]
    p.text = title.upper()
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.font.name = 'Arial'

    m_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9.33), Inches(2))
    tf = m_box.text_frame
    p = tf.paragraphs[0]
    p.text = main_text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)
    p.font.color.rgb = mute_color
    p.font.name = 'Arial'

    h_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.33), Inches(1.5))
    tf = h_box.text_frame
    p = tf.paragraphs[0]
    p.text = highlight_text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = accent_color
    p.font.name = 'Arial'
    
    return slide

# --- GENERATION ---

# Slide 1: Title
create_title_slide(
    prs, 
    "Quality Growth Value (QGV)", 
    "Fundamental Strategy",
    "Investing Where Risk is Managed and Uncertainty Creates Opportunity\n\nOleander Financial Technologies Pvt Ltd. | SEBI Registration No. INH000022491"
)

# Slide 2: Philosophy
create_content_slide(
    prs,
    "The Core Philosophy",
    "Redefining Risk and Returns",
    [
        "The Myth: High returns require taking excessive risk.",
        "The Reality: Sustainable, compounding returns are generated by minimizing risk while tactically navigating market uncertainty.",
        "Our Approach: We invest in fundamentally robust Mid & Small Cap businesses where temporary uncertainty creates valuation comfort—strictly avoiding any risk of permanent capital impairment."
    ],
    page=2
)

# Slide 3: What We Look For
create_content_slide(
    prs,
    "What We Look For",
    "The Non-Negotiable Pillars of QGV",
    [
        "• Strong & Scalable Business Model: Proven ability to grow.",
        "• Visible Earnings Growth: Clear line-of-sight to future cash flows.",
        "• Sensible Capital Allocation: Respect for shareholder capital.",
        "• Prudent Leverage: Clean balance sheets.",
        "• Reasonable Valuation: Distinct margin of safety at entry.",
        "• Re-rating Potential: Triggers driving institutional re-pricing."
    ],
    page=3
)

# Slide 4: The Selection Process
create_content_slide(
    prs,
    "The Selection Process",
    "A 4-Stage Mechanical Filtration Funnel",
    [
        "Step 1: The Quality Filter\n  Screening for ROCE consistency and balance sheet strength.",
        "Step 2: Growth Visibility\n  Identifying revenue scalability and structural sector tailwinds.",
        "Step 3: Valuation Comfort\n  Disciplined buying during market corrections or mispricings.",
        "Step 4: Risk Assessment\n  Deep-dive into corporate governance and leverage history."
    ],
    page=4
)

# Slide 5: Portfolio Construction
create_content_slide(
    prs,
    "Portfolio Construction",
    "How the Portfolio is Managed",
    [
        "Concentration: 15–20 High-Conviction Stocks.",
        "Market Cap Focus: Dedicated exclusively to Mid & Small Caps.",
        "Diversification: Sensibly spread across uncorrelated sectors.",
        "Active Management: Quarterly dynamic weighting adjustments.",
        "Discipline: Strictly avoiding overcrowded retail themes."
    ],
    page=5
)

# Slide 6: Why Mid & Small Caps?
create_content_slide(
    prs,
    "Why Mid & Small Caps?",
    "The Thesis for Generating Alpha",
    [
        "Greater Information Gaps: Less institutional coverage provides an edge.",
        "Higher Mispricing Probability: The market misjudges temporary headwinds.",
        "Earnings Acceleration: Smaller bases allow for explosive growth.",
        "Early-Stage Compounding: Buying the giants of tomorrow, today.",
        "\n\"Alpha emerges where uncertainty is misunderstood.\""
    ],
    page=6
)

# Slide 7: Risk Management
create_content_slide(
    prs,
    "Risk Management",
    "Minimizing Permanent Capital Loss by Actively Avoiding:",
    [
        "❌ Structurally Weak Businesses",
        "❌ Excessive Leverage or Debt Traps",
        "❌ Corporate Governance Red Flags",
        "❌ Pure Price-Momentum Plays without Fundamental Backing",
        "\nCoupled with strict Position Sizing Discipline to ensure no single failure derails the portfolio."
    ],
    page=7
)

# Slide 8: Strategy Positioning
create_content_slide(
    prs,
    "Strategy Positioning",
    "Setting the Correct Expectations",
    [
        "What We Are Not:",
        "❌ Not Deep Value (We don't buy value traps)",
        "❌ Not Pure Momentum (We require fundamental anchors)",
        "❌ Not Thematic (Not bound to a single sector narrative)",
        "\nWhat We Are:\nA highly disciplined intersection of Quality + Growth Visibility + Valuation Comfort."
    ],
    page=8
)

# Slide 9: Strategy At A Glance
create_content_slide(
    prs,
    "Strategy At A Glance",
    "Quick Facts & Details",
    [
        "Universe: Top 1500 Listed Companies",
        "Portfolio Size: 15–20 High-Conviction Stocks",
        "Rebalancing: Quarterly",
        "Core Style: Quality Growth at a Reasonable Valuation",
        "Pricing Options: ₹9,999/Yearly or ₹2,999/Quarterly (+ taxes)"
    ],
    page=9
)

# Slide 10: Conclusion
create_emphasis_slide(
    prs,
    "Quality Growth Value",
    "Investing in Businesses, Not Narratives.",
    "Contact Us: contact@olgo.in\nOleander Financial Technologies Pvt Ltd.",
    page=10
)

# Save
output_path = '/Users/shubhrakasana/Desktop/oleander/olgo_mobile/repos/Backtest/precision-backtest-engine/Quality_Growth_Value_Strategy.pptx'
prs.save(output_path)
print(f"Presentation saved to: {output_path}")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define colors
bg_color = RGBColor(15, 23, 42)  # Dark blue
accent_color = RGBColor(16, 185, 129)  # Green
text_color = RGBColor(248, 250, 252)  # White
secondary_text = RGBColor(148, 163, 184)  # Gray

def add_slide_with_title(prs, title, subtitle=None):
    """Add a slide with title and optional subtitle"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Set background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = accent_color
    
    # Add subtitle if provided
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(0.8))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.alignment = PP_ALIGN.CENTER
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = secondary_text
    
    return slide

def add_bullet_points(slide, bullets, left=0.5, top=3.5, width=4.5):
    """Add bullet points to a slide"""
    text_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(3.5))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = text_color
        p.space_after = Pt(12)
    
    return text_box

# Slide 1: Title Slide
slide1 = add_slide_with_title(prs, "The Champion Strategy", "Proprietary Alpha in Unpopular Sectors")
tagline_box = slide1.shapes.add_textbox(Inches(2), Inches(4), Inches(6), Inches(1))
tagline_frame = tagline_box.text_frame
tagline_frame.text = "Identifying Hidden Strength Where Others See Noise"
tagline_para = tagline_frame.paragraphs[0]
tagline_para.alignment = PP_ALIGN.CENTER
tagline_para.font.size = Pt(20)
tagline_para.font.color.rgb = text_color

# Slide 2: Core Philosophy
slide2 = add_slide_with_title(prs, "The Core Philosophy", "Identifying Strength Where Others See Noise")
bullets2 = [
    "The Contrarian Edge: Most participants chase 'hot' themes, creating overcrowding risk",
    "Our Advantage: Focus on unpopular, neglected sectors where value is mispriced",
    "The Precision: Proprietary models identify isolated pockets of hidden strength"
]
add_bullet_points(slide2, bullets2, left=0.5, top=3)

# Add image
try:
    slide2.shapes.add_picture(
        '/Users/shubhrakasana/.gemini/antigravity/brain/fb82cc32-205b-490d-bc5d-b39a6be90e7c/market_structure_accumulation_1771245796714.png',
        Inches(5.5), Inches(2.5), width=Inches(4)
    )
except:
    pass

# Slide 3: Alpha Source
slide3 = add_slide_with_title(prs, "The Alpha Source", "Capturing Structural Imbalances")
bullets3 = [
    "The Logic: Prices move due to Supply/Demand imbalances",
    "The Signal: Our model detects subtle shifts:\n  • Supply Exhaustion: Selling pressure dries up\n  • Smart Accumulation: Demand stabilizes",
    "The Result: Entry at inflection point for high-probability moves"
]
add_bullet_points(slide3, bullets3, left=1, top=3, width=8)

# Slide 4: Selection Process
slide4 = add_slide_with_title(prs, "The Selection Process", "A Rigorous, Multi-Layered Funnel")
bullets4 = [
    "Step 1: Proprietary Sector Scan\n  Identifying sectors with 'Accumulation Signatures'",
    "Step 2: Strength Confirmation\n  Verifying broad sector strength to minimize risk",
    "Step 3: Momentum Ranking\n  Buying the leaders of the new trend (RSNP Rank)"
]
add_bullet_points(slide4, bullets4, left=5, top=3, width=4.5)

# Add funnel image
try:
    slide4.shapes.add_picture(
        '/Users/shubhrakasana/.gemini/antigravity/brain/fb82cc32-205b-490d-bc5d-b39a6be90e7c/selection_funnel_1771245833605.png',
        Inches(0.5), Inches(2.5), width=Inches(4)
    )
except:
    pass

# Slide 5: Quality Overlay
slide5 = add_slide_with_title(prs, "Fundamentally Sound", "Strength with Safety")
bullets5 = [
    "Safety First: Strict Liquidity and Volatility filters. No junk.",
    "Structural Integrity: Price strength backed by valid market structure",
    "Disciplined Allocation:\n  • Max 15 Stocks (Concentrated)\n  • Max 3 per Industry (Diversified)"
]
add_bullet_points(slide5, bullets5, left=1.5, top=3, width=7)

# Slide 6: Performance
slide6 = add_slide_with_title(prs, "Validated by Results", "Delivering Alpha (2023-2026)")
bullets6 = [
    "CAGR: 28.87%\n  Consistent outperformance vs Benchmark (~15%)",
    "Sharpe Ratio: 1.17\n  Superior risk-adjusted returns",
    "Drawdown: -21.85%\n  Contrarian approach protects capital during corrections"
]
add_bullet_points(slide6, bullets6, left=0.5, top=3, width=4.5)

# Add performance chart
try:
    slide6.shapes.add_picture(
        '/Users/shubhrakasana/.gemini/antigravity/brain/fb82cc32-205b-490d-bc5d-b39a6be90e7c/performance_chart_concept_1771245860815.png',
        Inches(5.5), Inches(2.5), width=Inches(4)
    )
except:
    pass

# Slide 7: Dynamic Edge
slide7 = add_slide_with_title(prs, "The Dynamic Edge", "Intelligent Holding Periods")
content_box = slide7.shapes.add_textbox(Inches(2), Inches(3.5), Inches(6), Inches(2))
content_frame = content_box.text_frame
content_frame.text = "Beyond Fixed Timeframes\n\nWe hold winners as long as our Proprietary Signal remains valid"
for para in content_frame.paragraphs:
    para.alignment = PP_ALIGN.CENTER
    para.font.size = Pt(24)
    para.font.color.rgb = text_color

boost_box = slide7.shapes.add_textbox(Inches(2), Inches(5.5), Inches(6), Inches(1))
boost_frame = boost_box.text_frame
boost_frame.text = "Potential Return Boost: >31% CAGR"
boost_para = boost_frame.paragraphs[0]
boost_para.alignment = PP_ALIGN.CENTER
boost_para.font.size = Pt(32)
boost_para.font.bold = True
boost_para.font.color.rgb = accent_color

# Save presentation
prs.save('/Users/shubhrakasana/Desktop/oleander/olgo_mobile/repos/Backtest/precision-backtest-engine/Champion_Strategy_Presentation.pptx')
print("Presentation created successfully: Champion_Strategy_Presentation.pptx")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)  # Widescreen 16:9
prs.slide_height = Inches(7.5)

# Define colors
bg_color = RGBColor(15, 23, 42)    # Dark Slate (Background)
accent_color = RGBColor(16, 185, 129)  # Emerald Green (Accent)
text_color = RGBColor(248, 250, 252)   # White/Light Gray (Main Text)
mute_color = RGBColor(148, 163, 184)   # Slate Gray (Secondary Text)
darker_bg = RGBColor(30, 41, 59)       # Slightly lighter bg for shapes
highlight_color = RGBColor(255, 172, 65) # Orange for emphasis

def set_slide_background(slide):
    """Sets the dark background for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color

def add_footer(slide, page_num):
    """Adds a footer and page number."""
    # Footer Line
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.0), Inches(13.333), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = darker_bg
    shape.line.fill.background() # No border
    
    # Text
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(6), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "CONFIDENTIAL | The Contrarian Strength Strategy"
    p.font.size = Pt(10)
    p.font.color.rgb = mute_color
    
    # Page Num
    txBox = slide.shapes.add_textbox(Inches(12.3), Inches(7.1), Inches(1), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(page_num)
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(10)
    p.font.color.rgb = mute_color

def add_master_layout_elements(slide):
    """Adds common design elements like the accent bar."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.2), Inches(0.15), Inches(0.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()

def create_title_slide(prs, title, subtitle, tagline):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3), Inches(0.4), Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = accent_color
    shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.font.name = 'Arial'
    
    p = tf.add_paragraph()
    p.text = subtitle
    p.font.size = Pt(32)
    p.font.color.rgb = accent_color
    p.font.name = 'Arial'
    p.space_before = Pt(10)

    tag_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(1))
    tf = tag_box.text_frame
    p = tf.paragraphs[0]
    p.text = tagline
    p.font.size = Pt(20)
    p.font.color.rgb = mute_color
    p.font.name = 'Arial'

    return slide

def create_content_slide(prs, title, subtitle, bullets, image_path=None, page=0, disclaimer=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_master_layout_elements(slide)
    add_footer(slide, page)
    
    # Title
    t_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    tf = t_box.text_frame
    p = tf.paragraphs[0]
    p.text = title.upper()
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.font.name = 'Arial'
    
    p = tf.add_paragraph()
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = accent_color
    p.font.name = 'Arial'

    if image_path:
        text_width = Inches(6)
        img_left = Inches(7)
        img_top = Inches(2)
        img_width = Inches(5.8)
    else:
        text_width = Inches(10)
        img_left = 0
    
    # Bullets
    b_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), text_width, Inches(4.5))
    tf = b_box.text_frame
    tf.word_wrap = True
    
    for item in bullets:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(22)
        p.font.color.rgb = text_color
        p.font.name = 'Arial'
        p.space_after = Pt(20)
        p.space_before = Pt(5)
        p.level = 0
        
    if image_path:
        try:
            slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)
        except Exception as e:
            print(f"Could not load image {image_path}: {e}")
            pass
    
    if disclaimer:
        d_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.5))
        tf = d_box.text_frame
        p = tf.paragraphs[0]
        p.text = disclaimer
        p.font.size = Pt(12)
        p.font.color.rgb = mute_color
        p.font.italic = True

    return slide

def create_full_chart_slide(prs, title, subtitle, image_path, bullets=[], page=0, highlight_text=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_master_layout_elements(slide)
    add_footer(slide, page)
    
    t_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    tf = t_box.text_frame
    p = tf.paragraphs[0]
    p.text = title.upper()
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.font.name = 'Arial'
    
    p = tf.add_paragraph()
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = accent_color
    p.font.name = 'Arial'

    try:
        slide.shapes.add_picture(image_path, Inches(0.5), Inches(2), width=Inches(12.3))
    except Exception as e:
        print(f"Could not load image {image_path}: {e}")

    if highlight_text:
        h_box = slide.shapes.add_textbox(Inches(9.5), Inches(0.8), Inches(3.5), Inches(1))
        tf = h_box.text_frame
        p = tf.paragraphs[0]
        p.text = highlight_text
        p.alignment = PP_ALIGN.RIGHT
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = highlight_color
        
    return slide

def create_emphasis_slide(prs, title, main_text, highlight_text, page=0):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_footer(slide, page)
    
    t_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.33), Inches(1.5))
    tf = t_box.text_frame
    p = tf.paragraphs[0]
    p.text = title.upper()
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.font.name = 'Arial'

    m_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(9.33), Inches(2))
    tf = m_box.text_frame
    p = tf.paragraphs[0]
    p.text = main_text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)
    p.font.color.rgb = mute_color
    p.font.name = 'Arial'

    h_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.33), Inches(1.5))
    tf = h_box.text_frame
    p = tf.paragraphs[0]
    p.text = highlight_text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = accent_color
    p.font.name = 'Arial'
    
    return slide

# --- GENERATION ---

# Slide 1: Title (FINAL CHOICE)
create_title_slide(
    prs, 
    "The Contrarian Strength", 
    "Proprietary Alpha in Unpopular Sectors",
    "Finding Strength Where Others Look Away"
)

# Slide 2: Philosophy (Core)
create_content_slide(
    prs,
    "The Core Philosophy",
    "Identifying Strength Where Others See Noise",
    [
        "The Contrarian Edge: Most participants chase 'hot' themes, creating overcrowding risk.",
        "Our Advantage: Focus on ignored, neglected sectors where new momentum is emerging.", 
        "The Precision: We select only the strongest stocks in these 'ignored' industries."
    ],
    image_path='/Users/shubhrakasana/.gemini/antigravity/brain/fb82cc32-205b-490d-bc5d-b39a6be90e7c/market_structure_accumulation_1771245796714.png',
    page=2
)

# Slide 3: Alpha Source
create_content_slide(
    prs,
    "The Alpha Source",
    "Capturing Structural Imbalances",
    [
        "The Logic: Prices move due to Supply/Demand imbalances.",
        "The Signal: Our model detects subtle shifts in ignored sectors:",
        "• Supply Exhaustion: Selling pressure dries up",
        "• Smart Accumulation: Demand stabilizes",
        "The Result: Entry at inflection point for high-probability moves."
    ],
    page=3
)

# Slide 4: Selection Process
create_content_slide(
    prs,
    "The Selection Process",
    "A Rigorous, Multi-Layered Funnel",
    [
        "Step 1: Proprietary Sector Scan\nIdentifying sectors with 'Accumulation Signatures'",
        "Step 2: Strength Confirmation\nVerifying broad sector strength to minimize risk",
        "Step 3: Momentum Ranking\nBuying the leaders of the new trend (RSNP Rank)"
    ],
    image_path='/Users/shubhrakasana/.gemini/antigravity/brain/fb82cc32-205b-490d-bc5d-b39a6be90e7c/selection_funnel_1771245833605.png',
    page=4
)

# Slide 5: Quality Overlay
create_content_slide(
    prs,
    "Fundamentally Sound",
    "Strength with Safety",
    [
        "Safety First: Strict Liquidity and Volatility filters. No junk.",
        "Structural Integrity: Price strength backed by valid market structure.",
        "Disciplined Allocation:",
        "• Max 15 Stocks (Concentrated)",
        "• Max 3 per Industry (Diversified)"
    ],
    page=5
)

# Slide 6: Strategy Features 
create_content_slide(
    prs,
    "Strategy At A Glance",
    "Configuration & Structure",
    [
        "Benchmark: Nifty 500 / Top 1000 Universe",
        "Structure: Concentrated Portfolio (10-15 Stocks)",
        "Turnover: High Active Rotation (~7-10 stocks replaced/qtr)",
        "Pricing: 2% Fixed / 20% Profit Share (Above Hurdle)"
    ],
    image_path='/Users/shubhrakasana/.gemini/antigravity/brain/fb82cc32-205b-490d-bc5d-b39a6be90e7c/strategy_features_icons_1771299422555.png',
    page=6
)

# Slide 7: Long-Term Performance (9-Year Backtest)
create_content_slide(
    prs,
    "Long-Term Validation",
    "Full Cycle Backtest (2017-2026)",
    [
        "CAGR: 22.5% (vs 8.8% Bench)",
        "Sharpe (Rf=6%): 0.83 (vs 0.23 Bench)",
        "Max Drawdown: -41.1% (vs -64.0%)",
        "Total Return: 470% (vs 106%)"
    ],
    image_path='/Users/shubhrakasana/Desktop/oleander/olgo_mobile/repos/Backtest/precision-backtest-engine/outputs/full_nav_comparison.png',
    page=7,
    disclaimer="Note: We have back-tested the quantitative part of this strategy to give you an idea on how this thought process has performed in the past. Actual implementation includes qualitative fundamental checks."
)

# Slide 8: Yearly Returns breakdown
create_full_chart_slide(
    prs,
    "Annual Returns Breakdown",
    "Consistent Outperformance Across Cycles",
    image_path='/Users/shubhrakasana/Desktop/oleander/olgo_mobile/repos/Backtest/precision-backtest-engine/outputs/annual_returns_comparison.png',
    page=8,
    highlight_text="Yearly Outperformance: 80%"
)

# Slide 9: Quarterly Returns
create_full_chart_slide(
    prs,
    "Quarterly Performance",
    "Granular View of Alpha Generation",
    image_path='/Users/shubhrakasana/Desktop/oleander/olgo_mobile/repos/Backtest/precision-backtest-engine/outputs/quarterly_returns_comparison.png',
    page=9,
    highlight_text="Quarterly Outperformance: 75%"
)

# Slide 10: Dynamic Edge (Conclusion)
create_emphasis_slide(
    prs,
    "The Dynamic Edge",
    "We hold winners as long as our Proprietary Signal remains valid.",
    "Potential Return Boost: >31% CAGR",
    page=10
)

# Save
output_path = '/Users/shubhrakasana/Desktop/oleander/olgo_mobile/repos/Backtest/precision-backtest-engine/Contrarian_Strength_Presentation_v9.pptx'
prs.save(output_path)
print(f"Presentation v9 saved to: {output_path}")
